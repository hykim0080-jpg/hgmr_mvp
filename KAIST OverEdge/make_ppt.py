#!/usr/bin/env python3
"""
KAIST OverEdge 창업 아이디어 기술서 — 한글마루 PPT 생성
슬라이드 구성 (8장):
  1. 표지
  2. 요약 소개   (1p 제한 충족)
  3. Problem 1/2
  4. Problem 2/2  → 2p 제한 충족
  5. Solution 1/2
  6. Solution 2/2 → 2p 제한 충족
  7. AI 활용 역량 1/2
  8. AI 활용 역량 2/2 → 2p 제한 충족
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── 슬라이드 크기 ────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── 색상 팔레트 ──────────────────────────────────────
NAVY  = RGBColor(0x0D, 0x26, 0x44)
BLUE  = RGBColor(0x1A, 0x5F, 0x9E)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED   = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF4, 0xF6, 0xF8)
MGRAY = RGBColor(0xBD, 0xC3, 0xC7)
DGRAY = RGBColor(0x4A, 0x5A, 0x6A)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)

FONT = "Apple SD Gothic Neo"


# ── 헬퍼 함수 ────────────────────────────────────────

def slide_bg(sl, color):
    s = sl.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def rect(sl, l, t, w, h, fill=None):
    s = sl.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h, sz=12, bold=False,
        color=DGRAY, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name  = FONT
        r.font.size  = Pt(sz)
        r.font.bold  = bold
        r.font.italic = italic
        r.font.color.rgb = color

def set_cell(cell, text, sz=12, bold=False,
             fg=DGRAY, bg_hex=None, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = FONT
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.color.rgb = fg
    if bg_hex:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for tag in [qn('a:solidFill'), qn('a:noFill')]:
            for el in tcPr.findall(tag):
                tcPr.remove(el)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf,   qn('a:srgbClr'))
        sc.set('val', bg_hex.lstrip('#'))

def header(sl, title, sub=None, bg=NAVY, fh=Inches(1.15)):
    rect(sl, 0, 0, W, fh, fill=bg)
    txt(sl, title, Inches(0.45), Inches(0.18), W - Inches(0.9), Inches(0.65),
        sz=22, bold=True, color=WHITE)
    if sub:
        txt(sl, sub, Inches(0.45), Inches(0.72), W - Inches(0.9), Inches(0.38),
            sz=12, color=RGBColor(0xAD, 0xC6, 0xE5), italic=True)
    return fh


# ── 프레젠테이션 생성 ─────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, NAVY)

# 왼쪽 그린 바
rect(sl, 0, 0, Inches(0.45), H, fill=GREEN)

# 상단 레이블
txt(sl, "KAIST OverEdge  |  창업 아이디어 기술서",
    Inches(0.65), Inches(0.42), W - Inches(0.9), Inches(0.4),
    sz=13, color=RGBColor(0x7F, 0xB3, 0xE0))

# 메인 타이틀
txt(sl, "한글마루",
    Inches(0.65), Inches(1.3), Inches(8.5), Inches(1.7),
    sz=58, bold=True, color=WHITE)

# 그린 언더라인
rect(sl, Inches(0.65), Inches(2.95), Inches(4.0), Inches(0.09), fill=GREEN)

# 서브타이틀
txt(sl,
    "문맥 속 빈칸채우기 퀴즈를 통해\n게임처럼 즐겁게 문해력을 키우는\nAI 어휘 성장 플랫폼",
    Inches(0.65), Inches(3.1), Inches(9.0), Inches(2.3),
    sz=22, color=RGBColor(0xCC, 0xDC, 0xF0))

# 우측 MVP 현황 박스
rect(sl, Inches(10.2), Inches(1.9), Inches(2.88), Inches(3.7),
     fill=RGBColor(0x15, 0x3A, 0x60))
txt(sl, "MVP 현황", Inches(10.4), Inches(2.05), Inches(2.5), Inches(0.42),
    sz=13, bold=True, color=GREEN)
for i, s in enumerate([
    "✅  1,189개 큐레이션 단어",
    "✅  iOS · Android 배포",
    "✅  Claude API 자동화",
    "✅  Firebase 실시간 동기화",
    "✅  관리자 검수 시스템",
]):
    txt(sl, s, Inches(10.4), Inches(2.52) + i * Inches(0.62),
        Inches(2.5), Inches(0.52), sz=11, color=WHITE)

# 하단 바
rect(sl, 0, H - Inches(0.85), W, Inches(0.85), fill=RGBColor(0x08, 0x18, 0x2E))
txt(sl, "AI 어휘 성장 플랫폼  |  교육  |  2026  |  Hyeongi Kim",
    Inches(0.65), H - Inches(0.72), W - Inches(1), Inches(0.58),
    sz=11, color=MGRAY)


# ════════════════════════════════════════════════════
# SLIDE 2 — 요약 소개 (1p)
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, LGRAY)
hh = header(sl, "요약 소개", "Summary Overview  ·  전체 1 Page")

CW  = Inches(4.11)
CGP = Inches(0.15)
CT  = hh + Inches(0.28)
CH  = H - CT - Inches(0.28)

COLS = [
    (Inches(0.22),              RED,   "🔴  Problem",
     "풀고자 하는 문제",
     "AI 시대에 언어 구사력이 핵심 역량으로 부상한 반면,\n현대인의 어휘력은 숏폼 미디어로 급격히 저하.\n\n• 서울 고1 10명 중 3명 기초 문해력 미달\n• 생활 문해력 어려움 성인 735만 명\n• 외국어와 달리 모국어 어휘 전문 플랫폼 전무\n• KAIST에서도 어휘 소통 병목 직접 체감"),
    (Inches(0.22) + CW + CGP,  BLUE,  "🔵  Solution",
     "나의 솔루션",
     "한글마루: AI Agent 기반 한국어 어휘 성장 플랫폼\n\n• 자모 직접 타이핑 빈칸채우기 → 능동적 회상\n• 유저 간 실시간 오답 통계 → 메타인지 강화\n• 수능·학술·격식·감정별 큐레이션 세션\n• iOS · Android MVP 배포 완료"),
    (Inches(0.22) + (CW + CGP)*2, GREEN, "🟢  AI 활용 역량",
     "AI 도메인 전문성 및 활용 계획",
     "Claude API = AI 공동창업자\n1인 자동화 파이프라인 구축 완료\n\n• 수능 어휘 662개 예문 자동 생성\n• admin.html 검수 + 품질 자동 필터링\n• Firestore → 앱 실시간 동기화 스크립트\n• 현재 1,189개 큐레이션 단어 서비스 중"),
]

for (cl, color, label, sub, desc) in COLS:
    rect(sl, cl, CT, CW, CH, fill=color)
    txt(sl, label, cl+Inches(0.2), CT+Inches(0.14), CW-Inches(0.4), Inches(0.52),
        sz=15, bold=True, color=WHITE)
    rect(sl, cl, CT+Inches(0.7), CW, Inches(0.04), fill=WHITE)
    txt(sl, sub,  cl+Inches(0.2), CT+Inches(0.78), CW-Inches(0.4), Inches(0.38),
        sz=10, italic=True, color=RGBColor(0xD5, 0xEC, 0xFF))
    txt(sl, desc, cl+Inches(0.2), CT+Inches(1.22), CW-Inches(0.4), CH-Inches(1.32),
        sz=12, color=WHITE)


# ════════════════════════════════════════════════════
# SLIDE 3 — Problem 1/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "Problem  |  풀고자 하는 문제",
            "1 / 2  —  시대적 배경 · 시장 현황", RED)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=RED)

txt(sl, "① 시대적 배경 — AI가 언어력 격차를 심화시킨다",
    Inches(0.25), hh+Inches(0.18), W-Inches(0.4), Inches(0.4),
    sz=15, bold=True, color=RED)
txt(sl,
    "AI 보편화로 단순 실행력은 자동화되는 반면, 결과를 꿰뚫는 언어 구사력과 인문학적 사고력이 핵심 역량으로 재부상합니다. "
    "AI와 소통하는 프롬프트 작성에서도 의도를 정확히 언어화하는 능력은 기술적 기량이 됩니다. "
    "인공지능은 매일 수십억 개의 단어를 학습하며 똑똑해지는데, 우리는 '심심한 사과'의 뜻을 몰라 검색창을 헤매고 있습니다.",
    Inches(0.25), hh+Inches(0.64), W-Inches(0.4), Inches(0.85),
    sz=13, color=DGRAY)

rect(sl, Inches(0.25), hh+Inches(1.56), W-Inches(0.4), Inches(0.04), fill=MGRAY)

txt(sl, "② 시장 현황 — 어휘력 위기는 수치로 증명된다",
    Inches(0.25), hh+Inches(1.68), W-Inches(0.4), Inches(0.4),
    sz=15, bold=True, color=RED)

STATS = [
    ("30%",    "서울 고1 학생 중\n기초 문해력 미달\n(교육청 진단검사)"),
    ("25%",    "중2 학생 중\n수업 이해 불가 수준\n(서울시 교육청)"),
    ("735만 명", "생활 문해력에\n어려움 겪는 성인\n(한국교육개발원)"),
    ("KAIST\n체감", "이공계 최고 두뇌도\n어휘 부족으로\n소통 병목 경험"),
]
SW = Inches(3.0); SG = Inches(0.21); SY = hh+Inches(2.18); SH = Inches(2.6)
for i, (num, desc) in enumerate(STATS):
    sx = Inches(0.22) + i * (SW + SG)
    rect(sl, sx, SY, SW, SH, fill=RGBColor(0xFF, 0xEB, 0xEB))
    rect(sl, sx, SY, SW, Inches(0.08), fill=RED)
    txt(sl, num, sx, SY+Inches(0.15), SW, Inches(0.95),
        sz=26, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(sl, desc, sx, SY+Inches(1.1), SW, Inches(1.4),
        sz=12, color=DGRAY, align=PP_ALIGN.CENTER)

txt(sl, "출처: 서울시 교육청 진단검사 보고서, 한국교육개발원 성인문해능력조사",
    Inches(0.25), H-Inches(0.45), W-Inches(0.4), Inches(0.38),
    sz=10, color=MGRAY)


# ════════════════════════════════════════════════════
# SLIDE 4 — Problem 2/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "Problem  |  풀고자 하는 문제",
            "2 / 2  —  시장 공백 분석", RED)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=RED)

txt(sl, "③ 시장 공백 — 외국어는 넘치고, 모국어 어휘 전문 플랫폼은 전무하다",
    Inches(0.25), hh+Inches(0.18), W-Inches(0.4), Inches(0.4),
    sz=15, bold=True, color=RED)

TBL_DATA = [
    ["구분",       "외국어 학습 시장 (영어)",              "한국어 어휘 학습 (현황)"],
    ["대표 서비스", "Duolingo, Quizlet, 클래스101 등 수백 개",  "전무 (단순 검색 도구 수준)"],
    ["학습 방식",  "문맥 기반, 게이미피케이션, SRS 알고리즘",     "없음"],
    ["개인화",     "오답 분석, 맞춤형 복습 스케줄",               "없음"],
    ["시장 규모",  "국내 영어교육 시장 4조 원+",                  "미개척 블루오션"],
    ["잠재 수요",  "─",                                          "생활 문해력 어려움 성인 735만 명"],
]

tbl = sl.shapes.add_table(
    len(TBL_DATA), 3,
    Inches(0.22), hh+Inches(0.65),
    W-Inches(0.44), Inches(3.7)
).table
for i, w in enumerate([Inches(2.2), Inches(5.8), Inches(4.6)]):
    tbl.columns[i].width = w

BG_MAP = {
    (0, 0): '0D2644', (0, 1): '0D2644', (0, 2): '0D2644',
}
for ri, row in enumerate(TBL_DATA):
    for ci, text in enumerate(row):
        cell = tbl.cell(ri, ci)
        if ri == 0:
            set_cell(cell, text, sz=13, bold=True, fg=WHITE, bg_hex='0D2644',
                     align=PP_ALIGN.CENTER)
        elif ci == 0:
            set_cell(cell, text, sz=12, bold=True, fg=DARK,
                     bg_hex='F0F0F5' if ri%2==0 else 'E8E8F0')
        elif ci == 2:
            set_cell(cell, text, sz=12, bold=False, fg=RED,
                     bg_hex='FFF5F5' if ri%2==0 else 'FFFFFF')
        else:
            set_cell(cell, text, sz=12, bold=False, fg=DGRAY,
                     bg_hex='F4F6F8' if ri%2==0 else 'FFFFFF')

rect(sl, Inches(0.22), H-Inches(1.3), W-Inches(0.44), Inches(0.92),
     fill=RGBColor(0xFF, 0xEB, 0xEB))
rect(sl, Inches(0.22), H-Inches(1.3), Inches(0.07), Inches(0.92), fill=RED)
txt(sl,
    "💡  기회 요인: 모국어 어휘 학습 시장은 검증된 수요(735만 명)에도 불구하고 경쟁자가 없는 블루오션입니다.\n"
    "     AI Agent를 통해 1인이 고품질 콘텐츠를 대량 생산하는 것이 유일하고 확실한 진입 전략입니다.",
    Inches(0.38), H-Inches(1.24), W-Inches(0.65), Inches(0.84),
    sz=13, bold=True, color=RED)


# ════════════════════════════════════════════════════
# SLIDE 5 — Solution 1/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "Solution  |  정의한 문제에 대한 나의 솔루션",
            "1 / 2  —  3가지 차별화 솔루션", BLUE)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=BLUE)

SOL = [
    ("①",
     "원어민 맞춤형 능동적 학습  —  직접 타이핑하는 빈칸채우기",
     "이미 문장 구조를 이해하는 원어민에게 선택지 암기는 불필요합니다. 커스텀 한글 자모 키보드로 직접 타이핑하는 능동적 회상(Active Recall)으로 장기 기억 효율을 높입니다. "
     "정답을 맞힐 때마다 쌓이는 경험치(Lv.)와 연속 학습일 시스템으로, 지루한 국어 공부가 아닌 게임 같은 즐거운 루틴을 만듭니다."),
    ("②",
     "메타인지 강화  —  유저 간 실시간 오답 통계",
     "정답 제출 후 단순히 나의 정오답 여부만이 아니라 다른 유저들이 입력한 정오답 및 유의어 비율을 실시간으로 보여줍니다. "
     "'내가 빠진 함정에 남들도 똑같이 빠졌다'는 데이터를 통해 학습 좌절감 대신 함께 성장하는 강력한 연결감과 메타인지를 동시에 제공합니다."),
    ("③",
     "생애주기별 맞춤 큐레이션  —  지금 필요한 어휘 세션 직접 선택",
     "수능 고득점 목표 고등학생(필수·심화 수능 어휘 726개), 논문·발표 대학생(학술·논리 어휘), "
     "세대 간 소통 어려움 직장인(격식·비즈니스 어휘), 감정 표현을 원하는 성인(감정·심리 어휘). "
     "유저가 현재 가장 필요한 어휘를 직접 선택하여 즉각적인 효과를 제공합니다."),
]

IH = Inches(1.72); IG = Inches(0.1); IT = hh + Inches(0.22)
for i, (num, title, desc) in enumerate(SOL):
    iy = IT + i * (IH + IG)
    rect(sl, Inches(0.22), iy, Inches(0.6), IH, fill=BLUE)
    txt(sl, num, Inches(0.22), iy+Inches(0.58), Inches(0.6), Inches(0.58),
        sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bg_c = RGBColor(0xEE, 0xF5, 0xFF) if i % 2 == 0 else RGBColor(0xE4, 0xED, 0xF9)
    rect(sl, Inches(0.85), iy, W-Inches(1.1), IH, fill=bg_c)
    txt(sl, title, Inches(1.05), iy+Inches(0.1), W-Inches(1.3), Inches(0.48),
        sz=14, bold=True, color=BLUE)
    txt(sl, desc, Inches(1.05), iy+Inches(0.57), W-Inches(1.3), IH-Inches(0.64),
        sz=12, color=DGRAY)


# ════════════════════════════════════════════════════
# SLIDE 6 — Solution 2/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "Solution  |  정의한 문제에 대한 나의 솔루션",
            "2 / 2  —  큐레이션 세션 & 로드맵", BLUE)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=BLUE)

# 좌측: 큐레이션 테이블
txt(sl, "생애주기별 큐레이션 구성",
    Inches(0.25), hh+Inches(0.18), Inches(6.5), Inches(0.38),
    sz=14, bold=True, color=BLUE)

CUR = [
    ["타겟",              "큐레이션 세션",     "현재 단어 수"],
    ["수능 준비 고등학생", "필수·심화 수능 어휘", "726개 ✅"],
    ["논문·발표 대학생",  "학술·논리 어휘",     "85개 ✅"],
    ["사회초년생·직장인", "격식·비즈니스 어휘", "74개 ✅"],
    ["자기표현 원하는 성인", "감정·심리 어휘",  "준비 중 ⏳"],
    ["전체 통합",         "전 주제 세션",       "1,189개 ✅"],
]
ctbl = sl.shapes.add_table(
    len(CUR), 3, Inches(0.22), hh+Inches(0.62), Inches(6.7), Inches(3.5)
).table
for i, w in enumerate([Inches(2.45), Inches(2.75), Inches(1.5)]):
    ctbl.columns[i].width = w
for ri, row in enumerate(CUR):
    for ci, text in enumerate(row):
        cell = ctbl.cell(ri, ci)
        if ri == 0:
            set_cell(cell, text, sz=12, bold=True, fg=WHITE, bg_hex='1A5F9E',
                     align=PP_ALIGN.CENTER)
        elif ci == 2:
            c = GREEN if '✅' in text else MGRAY
            set_cell(cell, text, sz=12, bold=True, fg=c,
                     bg_hex='F0F5FF' if ri%2==0 else 'FFFFFF', align=PP_ALIGN.CENTER)
        else:
            set_cell(cell, text, sz=12, fg=DGRAY,
                     bg_hex='F0F5FF' if ri%2==0 else 'FFFFFF')

# 좌하단: 로드맵
rect(sl, Inches(0.22), H-Inches(1.15), Inches(6.7), Inches(0.85),
     fill=RGBColor(0xE8, 0xF0, 0xFF))
txt(sl,
    "로드맵:  [1단계] 앱 정식 런칭 → 유저 1,000명 확보\n"
    "           →  [2단계] SNS 바이럴 (어휘 해상도 테스트)\n"
    "           →  [3단계] 월간 구독형 프리미엄 BM 출시",
    Inches(0.38), H-Inches(1.08), Inches(6.3), Inches(0.8),
    sz=11, color=BLUE)

# 우측: MVP 현황
txt(sl, "MVP 구현 현황",
    Inches(7.1), hh+Inches(0.18), Inches(5.9), Inches(0.38),
    sz=14, bold=True, color=BLUE)

MVP_ITEMS = [
    ("✅", GREEN, "iOS / Android Capacitor 앱 배포"),
    ("✅", GREEN, "커스텀 한글 자모 입력 키보드"),
    ("✅", GREEN, "SRS 간격반복 알고리즘"),
    ("✅", GREEN, "Firebase 실시간 동기화"),
    ("✅", GREEN, "주제별 학습 세션 (65/35 가중치 적용)"),
    ("✅", GREEN, "유저 간 실시간 오답 통계"),
    ("✅", GREEN, "관리자 검수 시스템 (admin.html)"),
    ("⏳", MGRAY, "오답 복습 모드 (개발 예정)"),
    ("⏳", MGRAY, "AI 개인화 코치 (개발 예정)"),
]
for i, (icon, color, label) in enumerate(MVP_ITEMS):
    txt(sl, f"{icon}  {label}",
        Inches(7.1), hh+Inches(0.65) + i*Inches(0.42),
        Inches(5.9), Inches(0.38),
        sz=12, color=color)


# ════════════════════════════════════════════════════
# SLIDE 7 — AI 활용 역량 1/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "AI 활용 역량  |  AI 도메인 전문성 및 활용 계획",
            "1 / 2  —  현재 구현된 AI Agent 파이프라인", GREEN)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=GREEN)

txt(sl, "Claude API를 AI 공동창업자로  —  1인 풀스택 자동화 파이프라인 (현재 운영 중)",
    Inches(0.25), hh+Inches(0.15), W-Inches(0.4), Inches(0.42),
    sz=14, bold=True, color=GREEN)

# 파이프라인 다이어그램
PIPE = [
    ("Claude\nAPI",       GREEN),
    ("예문\n자동 생성",    RGBColor(0x16, 0x7A, 0x4B)),
    ("admin.html\n품질 검수", BLUE),
    ("Firestore\nDB 동기화", RGBColor(0xE8, 0x75, 0x22)),
    ("iOS · Android\n앱 배포", RED),
]
PW = Inches(2.18); PH = Inches(0.98); PG = Inches(0.18); AW = Inches(0.28)
PY = hh + Inches(0.7)
TOT = len(PIPE)*PW + (len(PIPE)-1)*(PG+AW)
PX0 = (W - TOT) / 2
for i, (lbl, col) in enumerate(PIPE):
    px = PX0 + i*(PW+PG+AW)
    rect(sl, px, PY, PW, PH, fill=col)
    txt(sl, lbl, px, PY+Inches(0.1), PW, PH-Inches(0.1),
        sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(PIPE)-1:
        txt(sl, "→", px+PW+PG, PY+Inches(0.18), AW, Inches(0.65),
            sz=20, bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

# 상세 테이블
AI_TBL = [
    ["AI 활용 영역",   "현재 구현 완료 내용",                              "성과"],
    ["콘텐츠 생성",    "Claude API로 수능 어휘 662개 예문 자동 작성\n(문맥 포함 예문, 빈칸 처리, 유의어 자동화)",
                                                                           "662개 생성·적용 완료"],
    ["품질 검수",      "관리자 검수 웹(admin.html) + AI 자동 필터 연동\n(word_reviews 기반 '수정 필요' 자동 제외)",
                                                                           '"수정 필요" 45개 자동 제외'],
    ["데이터 동기화",  "sync_words.js: Firestore → words.json 자동 반영\n(word_edits · word_deletions · word_reviews 통합)",
                                                                           "검수 즉시 앱 자동 반영"],
    ["운영 전체",      "1인이 기획·개발·AI 콘텐츠 생성·검수·배포 전 과정\n(Capacitor 모바일 앱 + Firebase + Claude API 통합)",
                                                                           "MVP 배포 완료\n1,189개 운영 중"],
]
at = sl.shapes.add_table(
    len(AI_TBL), 3,
    Inches(0.22), hh+Inches(1.88),
    W-Inches(0.44), Inches(3.85)
).table
for i, w in enumerate([Inches(2.2), Inches(7.5), Inches(2.9)]):
    at.columns[i].width = w
for ri, row in enumerate(AI_TBL):
    for ci, text in enumerate(row):
        cell = at.cell(ri, ci)
        if ri == 0:
            set_cell(cell, text, sz=13, bold=True, fg=WHITE, bg_hex='27AE60',
                     align=PP_ALIGN.CENTER)
        elif ci == 0:
            set_cell(cell, text, sz=12, bold=True, fg=DARK,
                     bg_hex='EDFAF3' if ri%2==0 else 'F9FFF9')
        elif ci == 2:
            set_cell(cell, text, sz=12, bold=True, fg=GREEN,
                     bg_hex='F0FBF4' if ri%2==0 else 'FFFFFF', align=PP_ALIGN.CENTER)
        else:
            set_cell(cell, text, sz=12, fg=DGRAY,
                     bg_hex='F4F9F6' if ri%2==0 else 'FFFFFF')


# ════════════════════════════════════════════════════
# SLIDE 8 — AI 활용 역량 2/2
# ════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, WHITE)
hh = header(sl, "AI 활용 역량  |  AI 도메인 전문성 및 활용 계획",
            "2 / 2  —  AI 경쟁 우위 & 향후 확장 로드맵", GREEN)
rect(sl, 0, hh, Inches(0.08), H - hh, fill=GREEN)

# 좌측: AI 경쟁 우위
txt(sl, "AI가 창출하는 3가지 핵심 경쟁 우위",
    Inches(0.25), hh+Inches(0.18), Inches(6.3), Inches(0.38),
    sz=14, bold=True, color=GREEN)

ADV = [
    ("비용 구조 혁신",
     "단어 콘텐츠 생산에 전통적으로 국어교사·편집자 팀이 필요.\n"
     "→  Claude API 자동화로 1인이 1,189개 고품질 데이터 구축 완료.\n"
     "→  경쟁자 대비 콘텐츠 생산 비용 90% 이상 절감."),
    ("속도 우위",
     "새로운 주제 단어팩(의학·법률·반도체) 추가 시:\n"
     "기존 방식: 수개월 소요  →  AI 파이프라인: 수일 내 완료.\n"
     "→  시장 대응 속도 자체가 진입 장벽."),
    ("글로벌 확장성",
     "동일 파이프라인으로 영어·일본어 어휘팩 추가 시 글로벌 진출 즉시 가능.\n"
     "→  1인 창업자가 다국어 플랫폼을 운영할 수 있는 유일한 방법.\n"
     "→  AI 없이는 구조적으로 불가능한 비즈니스 모델."),
]
for i, (title, desc) in enumerate(ADV):
    ay = hh + Inches(0.65) + i * Inches(1.62)
    rect(sl, Inches(0.22), ay, Inches(6.3), Inches(1.52), fill=RGBColor(0xF0, 0xFB, 0xF4))
    rect(sl, Inches(0.22), ay, Inches(0.08), Inches(1.52), fill=GREEN)
    txt(sl, title, Inches(0.42), ay+Inches(0.1), Inches(5.9), Inches(0.42),
        sz=13, bold=True, color=GREEN)
    txt(sl, desc, Inches(0.42), ay+Inches(0.55), Inches(5.9), Inches(0.92),
        sz=11, color=DGRAY)

# 우측: 향후 로드맵
txt(sl, "향후 AI Agent 확장 로드맵",
    Inches(6.8), hh+Inches(0.18), Inches(6.2), Inches(0.38),
    sz=14, bold=True, color=GREEN)

ROAD = [
    ("단기  2026 Q3", GREEN, [
        "학습자 오답 패턴 AI 분석 → 개인화 복습 세션 자동 생성",
        "망각 곡선 연동 AI 푸시 알림",
        "단어팩 3,000개 이상으로 확장 (Claude API 자동화)",
    ]),
    ("중기  2026 Q4", BLUE, [
        "AI 학습 코치: 오늘 이 단어를 복습해야 하는 이유 설명",
        "어휘 해상도 테스트 SNS 배포 → 바이럴 유입 확대",
        "월간 구독형 프리미엄 모델 정식 출시",
    ]),
    ("장기  2027+", RED, [
        "영어·일본어 어휘팩 → 글로벌 시장 진출",
        "B2B: 학원·기업 어휘 교육 AI 솔루션",
        "AI 튜터 기반 완전 개인화 학습 플랫폼",
    ]),
]
rmy = hh + Inches(0.62)
for phase, col, items in ROAD:
    rect(sl, Inches(6.8), rmy, Inches(6.17), Inches(1.72), fill=LGRAY)
    rect(sl, Inches(6.8), rmy, Inches(0.25), Inches(1.72), fill=col)
    txt(sl, phase, Inches(7.15), rmy+Inches(0.12), Inches(5.8), Inches(0.42),
        sz=12, bold=True, color=col)
    for j, item in enumerate(items):
        txt(sl, f"• {item}",
            Inches(7.15), rmy+Inches(0.58) + j*Inches(0.36),
            Inches(5.7), Inches(0.35),
            sz=11, color=DGRAY)
    rmy += Inches(1.82)


# ── 저장 ─────────────────────────────────────────────
OUT = "/Users/hyk/Desktop/hgmr/KAIST OverEdge/KAIST_OverEdge_한글마루.pptx"
prs.save(OUT)
print(f"✅  저장 완료: {OUT}")
print(f"   슬라이드 수: {len(prs.slides)}장")
